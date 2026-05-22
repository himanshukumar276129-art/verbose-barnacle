import os
import json
import uuid
import httpx
from sqlmodel import Session
from fastapi import HTTPException
from pptx import Presentation
from docx import Document
from openpyxl import Workbook
from app.core.config import settings


class DocumentService:
    """Service to programmatically generate PPT, Word, and Excel files using multiple LLM providers."""

    @staticmethod
    def _clean_json_response(text: str) -> str:
        """Strip markdown markers if returned by the LLM."""
        text = text.strip()
        if text.startswith("```"):
            lines = text.split("\n")
            if lines[0].startswith("```"):
                lines = lines[1:]
            if lines[-1].startswith("```"):
                lines = lines[:-1]
            text = "\n".join(lines).strip()
        return text

    @staticmethod
    def _parse_json_safe(cleaned: str) -> dict:
        """Parse JSON from cleaned LLM output with fallback extraction."""
        try:
            return json.loads(cleaned)
        except Exception as e:
            # Fallback parsing helper if JSON starts/ends with extra text
            try:
                start = cleaned.find("{")
                end = cleaned.rfind("}")
                if start != -1 and end != -1:
                    return json.loads(cleaned[start:end+1])
                start_arr = cleaned.find("[")
                end_arr = cleaned.rfind("]")
                if start_arr != -1 and end_arr != -1:
                    return json.loads(cleaned[start_arr:end_arr+1])
            except Exception:
                pass
            raise HTTPException(status_code=500, detail=f"Failed to parse structured JSON from compiler: {e}")

    # =========================================================================
    #  Provider-specific content generators
    # =========================================================================

    @staticmethod
    async def _generate_via_groq(prompt: str, system_prompt: str, tier: int = 1) -> dict:
        """Generate structured JSON content using Groq (llama-3.3-70b-versatile)."""
        from app.services.providers.groq_provider import GroqProvider

        result = await GroqProvider.run_model(
            "llama-3.3-70b-versatile",
            {
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                "max_tokens": 4096,
                "temperature": 0.7
            },
            tier
        )

        # Extract content from OpenAI-compatible response
        if isinstance(result, dict) and "choices" in result:
            content = result["choices"][0]["message"]["content"]
        elif isinstance(result, str):
            content = result
        else:
            raise HTTPException(status_code=500, detail=f"Unexpected Groq response format: {result}")

        cleaned = DocumentService._clean_json_response(content)
        return DocumentService._parse_json_safe(cleaned)

    @staticmethod
    async def _generate_via_ollama(prompt: str, system_prompt: str, tier: int = 1) -> dict:
        """Generate structured JSON content using Ollama (llama3)."""
        from app.services.providers.ollama_provider import OllamaProvider

        result = await OllamaProvider.run_model(
            "llama3",
            {
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.7
            },
            tier
        )

        # Extract content from OpenAI-compatible response
        if isinstance(result, dict) and "choices" in result:
            content = result["choices"][0]["message"]["content"]
        elif isinstance(result, str):
            content = result
        else:
            raise HTTPException(status_code=500, detail=f"Unexpected Ollama response format: {result}")

        cleaned = DocumentService._clean_json_response(content)
        return DocumentService._parse_json_safe(cleaned)

    @staticmethod
    async def _generate_via_compiler(prompt: str, system_prompt: str) -> dict:
        """Generate structured JSON content using the dedicated DOCUMENT_COMPILER_KEY (Gemini or OpenAI)."""
        api_key = settings.DOCUMENT_COMPILER_KEY
        if not api_key:
            raise HTTPException(
                status_code=501,
                detail="Document compilation feature is currently inactive. Please configure DOCUMENT_COMPILER_KEY in your .env file to activate."
            )

        api_key = api_key.strip()
        
        # 1. If key starts with sk-, treat as OpenAI API key
        if api_key.startswith("sk-"):
            url = "https://api.openai.com/v1/chat/completions"
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            payload = {
                "model": "gpt-4o-mini",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                "response_format": {"type": "json_object"}
            }
            try:
                async with httpx.AsyncClient(timeout=60.0) as client:
                    response = await client.post(url, headers=headers, json=payload)
                    if response.status_code != 200:
                        raise Exception(f"OpenAI error response: {response.text}")
                    res_data = response.json()
                    content = res_data["choices"][0]["message"]["content"]
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"OpenAI compilation failed: {e}")

        # 2. Otherwise, treat as Google Gemini API Key
        else:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}"
            payload = {
                "contents": [{
                    "parts": [
                        {"text": f"System Guidelines: {system_prompt}"},
                        {"text": f"User Prompt: {prompt}"}
                    ]
                }],
                "generationConfig": {
                    "responseMimeType": "application/json"
                }
            }
            try:
                async with httpx.AsyncClient(timeout=60.0) as client:
                    response = await client.post(url, json=payload)
                    if response.status_code != 200:
                        raise Exception(f"Gemini error response: {response.text}")
                    res_data = response.json()
                    content = res_data["candidates"][0]["content"]["parts"][0]["text"]
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Gemini compilation failed: {e}")

        cleaned = DocumentService._clean_json_response(content)
        return DocumentService._parse_json_safe(cleaned)

    @staticmethod
    async def _generate_content(prompt: str, system_prompt: str, provider: str = "document_compiler", tier: int = 1) -> dict:
        """Route to the correct LLM backend based on provider."""
        if provider == "groq":
            return await DocumentService._generate_via_groq(prompt, system_prompt, tier)
        elif provider == "ollama":
            return await DocumentService._generate_via_ollama(prompt, system_prompt, tier)
        else:
            return await DocumentService._generate_via_compiler(prompt, system_prompt)

    # =========================================================================
    #  Document Generators (PPT, Word, Excel)
    # =========================================================================

    @staticmethod
    async def generate_ppt(prompt: str, base_url: str, provider: str = "document_compiler", tier: int = 1) -> str:
        """Generate a PPT presentation based on prompt."""
        system_prompt = (
            "You are a professional presentation planner. Output ONLY a valid JSON object with a single key 'slides' containing an array of objects. "
            "Do not include any markdown format, backticks, or text before/after the JSON.\n"
            "Format:\n"
            "{\n"
            "  \"slides\": [\n"
            "    {\n"
            "      \"title\": \"Slide Title\",\n"
            "      \"bullet_points\": [\"Point 1\", \"Point 2\", \"Point 3\"]\n"
            "    }\n"
            "  ]\n"
            "}"
        )
        
        data = await DocumentService._generate_content(prompt, system_prompt, provider, tier)
        slides = data.get("slides", [])
        if not isinstance(slides, list):
            slides = [data]

        prs = Presentation()
        # Set to standard wide 16:9 aspect ratio
        prs.slide_width = 12192000
        prs.slide_height = 6858000

        for slide_data in slides:
            # 1 is Title + Content layout
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Populate title
            slide.shapes.title.text = slide_data.get("title", "Untitled Slide")
            
            # Populate body/bullets
            placeholders = slide.placeholders
            if len(placeholders) > 1:
                tf = placeholders[1].text_frame
                tf.text = ""
                for i, pt in enumerate(slide_data.get("bullet_points", [])):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = pt
                    p.level = 0

        # Save presentation
        os.makedirs("static/generated", exist_ok=True)
        filename = f"ppt_{uuid.uuid4().hex[:12]}.pptx"
        file_path = os.path.join("static/generated", filename)
        prs.save(file_path)
        
        return f"{base_url}static/generated/{filename}"

    @staticmethod
    async def generate_word(prompt: str, base_url: str, provider: str = "document_compiler", tier: int = 1) -> str:
        """Generate a Word document based on prompt."""
        system_prompt = (
            "You are a professional document writer. Output ONLY a valid JSON object representing the document outline. "
            "Do not include any markdown format, backticks, or text before/after the JSON. "
            "Format:\n"
            "{\n"
            "  \"title\": \"Document Title\",\n"
            "  \"sections\": [\n"
            "    {\n"
            "      \"heading\": \"Section Title\",\n"
            "      \"paragraphs\": [\"First paragraph text.\", \"Second paragraph text.\"]\n"
            "    }\n"
            "  ]\n"
            "}"
        )
        
        data = await DocumentService._generate_content(prompt, system_prompt, provider, tier)
        
        doc = Document()
        doc.add_heading(data.get("title", "Document"), level=0)
        
        for section in data.get("sections", []):
            doc.add_heading(section.get("heading", ""), level=1)
            for p_text in section.get("paragraphs", []):
                doc.add_paragraph(p_text)

        # Save document
        os.makedirs("static/generated", exist_ok=True)
        filename = f"doc_{uuid.uuid4().hex[:12]}.docx"
        file_path = os.path.join("static/generated", filename)
        doc.save(file_path)
        
        return f"{base_url}static/generated/{filename}"

    @staticmethod
    async def generate_excel(prompt: str, base_url: str, provider: str = "document_compiler", tier: int = 1) -> str:
        """Generate an Excel spreadsheet based on prompt."""
        system_prompt = (
            "You are an expert database analyst. Output ONLY a valid JSON object representing a data spreadsheet. "
            "Do not include any markdown format, backticks, or text before/after the JSON. "
            "Format:\n"
            "{\n"
            "  \"sheet_title\": \"Sheet Name\",\n"
            "  \"headers\": [\"Column 1\", \"Column 2\", \"Column 3\"],\n"
            "  \"rows\": [\n"
            "    [\"Value 1a\", \"Value 1b\", \"Value 1c\"],\n"
            "    [\"Value 2a\", \"Value 2b\", \"Value 2c\"]\n"
            "  ]\n"
            "}"
        )
        
        data = await DocumentService._generate_content(prompt, system_prompt, provider, tier)
        
        wb = Workbook()
        ws = wb.active
        ws.title = data.get("sheet_title", "DataSheet")[:30] # Excel limit is 31 chars
        
        # Append headers
        ws.append(data.get("headers", []))
        
        # Append rows
        for row in data.get("rows", []):
            ws.append(row)

        # Style header row (optional but nice)
        for col_num in range(1, len(data.get("headers", [])) + 1):
            cell = ws.cell(row=1, column=col_num)
            cell.font = cell.font.copy(bold=True)

        # Save spreadsheet
        os.makedirs("static/generated", exist_ok=True)
        filename = f"excel_{uuid.uuid4().hex[:12]}.xlsx"
        file_path = os.path.join("static/generated", filename)
        wb.save(file_path)
        
        return f"{base_url}static/generated/{filename}"

    # =========================================================================
    #  PDF Document Generator (fpdf2 – Professional Indigo Theme)
    # =========================================================================

    @staticmethod
    async def generate_pdf(prompt: str, base_url: str, provider: str = "document_compiler", tier: int = 1) -> str:
        """Generate a professionally styled PDF document based on prompt using AI-structured content."""
        from fpdf import FPDF
        from datetime import datetime as _dt

        system_prompt = (
            "You are a professional document writer. Output ONLY a valid JSON object representing a document outline. "
            "Do not include any markdown format, backticks, or text before/after the JSON.\n"
            "Format:\n"
            "{\n"
            "  \"title\": \"Document Title\",\n"
            "  \"subtitle\": \"A brief one-line description\",\n"
            "  \"sections\": [\n"
            "    {\n"
            "      \"heading\": \"Section Title\",\n"
            "      \"paragraphs\": [\"First paragraph text.\", \"Second paragraph text.\"]\n"
            "    }\n"
            "  ]\n"
            "}"
        )

        data = await DocumentService._generate_content(prompt, system_prompt, provider, tier)

        # ── Theme Constants ──────────────────────────────────────────
        INDIGO = (79, 70, 229)        # Primary accent
        CHARCOAL = (31, 41, 55)       # Body text
        LIGHT_GRAY = (243, 244, 246)  # Subtle background fills
        WHITE = (255, 255, 255)
        DARK_GRAY = (107, 114, 128)   # Footer / meta text

        class VedaPDF(FPDF):
            """Custom FPDF subclass that draws themed headers and page-numbered footers."""
            def __init__(self, doc_title: str):
                super().__init__()
                self.doc_title = doc_title

            def header(self):
                # Indigo header band
                self.set_fill_color(*INDIGO)
                self.rect(0, 0, self.w, 18, 'F')
                self.set_font("Helvetica", "B", 11)
                self.set_text_color(*WHITE)
                self.set_xy(self.l_margin, 4)
                self.cell(self.epw, 10, self.doc_title[:80], align="L")
                
                # Timestamp on the right
                self.set_font("Helvetica", "", 8)
                self.set_xy(self.w - self.r_margin - 50, 5)
                self.cell(50, 8, _dt.utcnow().strftime("%Y-%m-%d %H:%M UTC"), align="R")
                self.set_y(22)

            def footer(self):
                self.set_y(-15)
                self.set_font("Helvetica", "I", 8)
                self.set_text_color(*DARK_GRAY)
                self.cell(0, 10, f"Page {self.page_no()}/{{nb}}", align="C")

        doc_title = data.get("title", "Document")
        doc_subtitle = data.get("subtitle", "")
        sections = data.get("sections", [])
        if not isinstance(sections, list):
            sections = [data]

        pdf = VedaPDF(doc_title)
        pdf.alias_nb_pages()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()

        # ── Cover Title Block ────────────────────────────────────────
        pdf.set_y(30)
        pdf.set_font("Helvetica", "B", 26)
        pdf.set_text_color(*INDIGO)
        pdf.multi_cell(pdf.epw, 12, doc_title, align="C")
        pdf.ln(4)

        if doc_subtitle:
            pdf.set_font("Helvetica", "", 12)
            pdf.set_text_color(*DARK_GRAY)
            pdf.multi_cell(pdf.epw, 8, doc_subtitle, align="C")
            pdf.ln(2)

        # Accent divider line
        pdf.ln(6)
        y_line = pdf.get_y()
        pdf.set_draw_color(*INDIGO)
        pdf.set_line_width(0.8)
        pdf.line(pdf.l_margin + 20, y_line, pdf.w - pdf.r_margin - 20, y_line)
        pdf.ln(10)

        # ── Body Sections ────────────────────────────────────────────
        for idx, section in enumerate(sections):
            heading = section.get("heading", f"Section {idx + 1}")
            paragraphs = section.get("paragraphs", [])

            # Section heading with indigo left-bar accent
            y_before = pdf.get_y()
            if y_before > 250: # Slightly more conservative page break
                pdf.add_page()
                y_before = pdf.get_y()

            pdf.set_fill_color(*INDIGO)
            pdf.rect(pdf.l_margin, y_before, 2.5, 9, 'F')
            pdf.set_xy(pdf.l_margin + 6, y_before)
            pdf.set_font("Helvetica", "B", 14)
            pdf.set_text_color(*CHARCOAL)
            pdf.cell(pdf.epw - 6, 9, heading)
            pdf.ln(12)

            # Paragraphs
            pdf.set_font("Helvetica", "", 11)
            pdf.set_text_color(*CHARCOAL)
            for para in paragraphs:
                if not isinstance(para, str):
                    para = str(para)
                pdf.multi_cell(pdf.epw, 6.5, para)
                pdf.ln(3)

            pdf.ln(4)

        # ── Save PDF ─────────────────────────────────────────────────
        os.makedirs("static/generated", exist_ok=True)
        filename = f"pdf_{uuid.uuid4().hex[:12]}.pdf"
        file_path = os.path.join("static/generated", filename)
        pdf.output(file_path)

        return f"{base_url}static/generated/{filename}"


